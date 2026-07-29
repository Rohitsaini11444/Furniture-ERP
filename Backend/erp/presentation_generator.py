import os
from decimal import Decimal
from io import BytesIO
from django.utils import timezone
from django.conf import settings
from PIL import Image as PILImage

# python-pptx imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


def find_image_path(item):
    """
    Finds absolute image file path for a Sample or BuyerMaster item.
    Checks model fields, related models (SampleImage, BuyerMasterFinishingImage),
    associated models, and media storage folders.
    """
    # 1. Direct image field on sample or buyer master
    for field_name in ['image', 'packaging_image']:
        img_attr = getattr(item, field_name, None)
        if img_attr and hasattr(img_attr, 'path') and os.path.exists(img_attr.path):
            return img_attr.path

    # 2. Check SampleImage related objects (item.images) if item is a Sample
    if hasattr(item, 'images'):
        try:
            first_img = item.images.filter(image__isnull=False).exclude(image='').first()
            if first_img and hasattr(first_img.image, 'path') and os.path.exists(first_img.image.path):
                return first_img.image.path
        except Exception:
            pass

    # 3. If item is a BuyerMaster, check associated Sample and its images
    sample = getattr(item, 'sample', None)
    if sample:
        if sample.image and hasattr(sample.image, 'path') and os.path.exists(sample.image.path):
            return sample.image.path
        if hasattr(sample, 'images'):
            try:
                first_img = sample.images.filter(image__isnull=False).exclude(image='').first()
                if first_img and hasattr(first_img.image, 'path') and os.path.exists(first_img.image.path):
                    return first_img.image.path
            except Exception:
                pass

    # 4. Check finishing_images if item is BuyerMaster
    if hasattr(item, 'finishing_images'):
        try:
            first_finish = item.finishing_images.filter(image__isnull=False).exclude(image='').first()
            if first_finish and hasattr(first_finish.image, 'path') and os.path.exists(first_finish.image.path):
                return first_finish.image.path
        except Exception:
            pass

    # 5. Check associated BuyerMaster packaging_image if item is Sample
    if hasattr(item, 'buyer_masters'):
        try:
            bm = item.buyer_masters.filter(packaging_image__isnull=False).exclude(packaging_image='').first()
            if bm and hasattr(bm.packaging_image, 'path') and os.path.exists(bm.packaging_image.path):
                return bm.packaging_image.path
        except Exception:
            pass

    # 6. Check media directory for matching file name by style_no or sample_id
    media_root = getattr(settings, 'MEDIA_ROOT', os.path.join(os.path.dirname(__file__), '..', 'media'))
    identifiers = [getattr(item, 'style_no', ''), getattr(item, 'sample_id', '')]
    identifiers = [i for i in identifiers if i]

    for identifier in identifiers:
        for ext in ['.jpg', '.jpeg', '.png', '.webp']:
            for folder in ['samples', 'buyer_masters', 'buyer_masters/packaging', 'buyer_masters/finishing', '']:
                test_path = os.path.join(media_root, folder, f"{identifier}{ext}")
                if os.path.exists(test_path):
                    return test_path

    return None


def generate_pptx_presentation(buyer, items, items_per_slide=2, include_price=True, include_specs=True):
    """
    Generates a 16:9 widescreen PowerPoint Presentation (.pptx)
    - Slide 1: Cover Page with Company Logo, Title, Prepared For Buyer, Date.
    - Slide 2..N: Dynamic Selected Items per Slide (Left/Right side-by-side: Image + Specs Table).
    - Slide N+1: Clean, Centered Thank You & Contact Info Card Slide.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    C_WALNUT = RGBColor(139, 90, 43)      # #8B5A2B Brand Theme
    C_GOLD = RGBColor(217, 119, 6)        # #D97706 Gold Accent
    C_CREAM_BG = RGBColor(248, 246, 242)  # #F8F6F2 Light Luxury Background
    C_DARK = RGBColor(30, 41, 59)         # #1E293B Primary Text
    C_SLATE = RGBColor(71, 85, 105)       # #475569 Subtitle Text
    C_MUTED = RGBColor(100, 116, 139)     # #64748B Label Muted Text
    C_WHITE = RGBColor(255, 255, 255)
    C_ROW_ALT = RGBColor(248, 250, 252)   # #F8FAFC Table Alt Row
    C_BORDER = RGBColor(226, 232, 240)    # #E2E8F0 Table Border

    # ── SLIDE 1: Cover Page ──
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_CREAM_BG
    bg1.line.fill.background()

    # Brand Header Bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.3))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_WALNUT
    top_bar.line.fill.background()

    # Company Logo Insertion
    logo_path = os.path.join(settings.BASE_DIR, '..', 'Frontend', 'src', 'assets', 'Pinkcity_Logo.png')
    if not os.path.exists(logo_path):
        logo_path = r"C:\Users\User\OneDrive\Desktop\ERP Furniture\Frontend\src\assets\Pinkcity_Logo.png"

    text_left = Inches(0.4)
    if os.path.exists(logo_path):
        try:
            slide1.shapes.add_picture(logo_path, Inches(0.4), Inches(0.15), width=Inches(1.0), height=Inches(1.0))
            text_left = Inches(1.6)
        except Exception as e:
            print("Error rendering logo in PPT:", e)

    tf_brand = slide1.shapes.add_textbox(text_left, Inches(0.25), Inches(11.0), Inches(0.8)).text_frame
    tf_brand.word_wrap = True
    p_brand = tf_brand.paragraphs[0]
    p_brand.text = "PINKCITY ENTERPRISES"
    p_brand.font.size = Pt(26)
    p_brand.font.bold = True
    p_brand.font.color.rgb = C_WHITE

    # Accent Gold Line
    gold_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), Inches(13.333), Inches(0.06))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = C_GOLD
    gold_line.line.fill.background()

    # Title & Subtitle Card
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.333), Inches(2.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_title = tf_title.paragraphs[0]
    p_title.text = "EXCLUSIVE FURNITURE COLLECTION"
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = C_DARK

    p_sub = tf_title.add_paragraph()
    p_sub.text = "Sample & Product Catalog Presentation"
    p_sub.font.size = Pt(20)
    p_sub.font.bold = True
    p_sub.font.color.rgb = C_WALNUT

    # Prepared For Card Box
    card_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.3), Inches(11.333), Inches(2.5))
    card_box.fill.solid()
    card_box.fill.fore_color.rgb = C_WHITE
    card_box.line.color.rgb = C_BORDER

    tf_info = card_box.text_frame
    tf_info.word_wrap = True
    
    p_b1 = tf_info.paragraphs[0]
    p_b1.text = "PREPARED SPECIALLY FOR:"
    p_b1.font.size = Pt(12)
    p_b1.font.bold = True
    p_b1.font.color.rgb = C_MUTED

    p_b2 = tf_info.add_paragraph()
    p_b2.text = f"{buyer.name} ({buyer.code})" if buyer else "Valued Client / General Presentation"
    p_b2.font.size = Pt(22)
    p_b2.font.bold = True
    p_b2.font.color.rgb = C_DARK

    if buyer and buyer.email:
        p_b3 = tf_info.add_paragraph()
        p_b3.text = f"Email: {buyer.email}"
        p_b3.font.size = Pt(14)
        p_b3.font.color.rgb = C_SLATE

    p_dt = tf_info.add_paragraph()
    p_dt.text = f"Date: {timezone.now().strftime('%d %B %Y')}"
    p_dt.font.size = Pt(14)
    p_dt.font.color.rgb = C_SLATE

    # ── SLIDES 2..N: Dynamic Selected Items per Slide ──
    chunk_size = max(1, min(2, items_per_slide))
    item_chunks = [items[i:i + chunk_size] for i in range(0, len(items), chunk_size)]

    for page_idx, chunk in enumerate(item_chunks, start=1):
        slide = prs.slides.add_slide(blank_layout)
        
        # Background
        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = RGBColor(250, 250, 249)
        s_bg.line.fill.background()

        # Header Bar Banner
        h_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.3), Inches(12.533), Inches(0.7))
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = C_WALNUT
        h_bar.line.fill.background()
        
        h_tf = h_bar.text_frame
        h_tf.word_wrap = True
        h_p = h_tf.paragraphs[0]
        h_p.text = f"  PINKCITY ENTERPRISES   |   PRODUCT CATALOG (Page {page_idx} of {len(item_chunks)})"
        h_p.font.size = Pt(16)
        h_p.font.bold = True
        h_p.font.color.rgb = C_WHITE
        h_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if chunk_size == 1:
            # ── 1 ITEM PER SLIDE: FULL PAGE LAYOUT (Spans 100% Slide Width) ──
            item = chunk[0]
            global_idx = page_idx
            item_left = Inches(0.4)
            full_width = Inches(12.533)

            style_val = getattr(item, 'style_no', '') or getattr(item, 'sample_id', '') or f"Item #{global_idx}"
            prod_val = getattr(item, 'product_name', 'Furniture Style')

            # Full Width Item Section Title Bar
            item_title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, item_left, Inches(1.15), full_width, Inches(0.5))
            item_title_box.fill.solid()
            item_title_box.fill.fore_color.rgb = C_WALNUT
            item_title_box.line.fill.background()
            
            t_tf = item_title_box.text_frame
            t_tf.word_wrap = True
            t_p = t_tf.paragraphs[0]
            t_p.text = f"  {global_idx}. {prod_val}   |   Style No: {style_val}"
            t_p.font.size = Pt(13)
            t_p.font.bold = True
            t_p.font.color.rgb = C_WHITE
            item_title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Left Box: Product Image Container (Width 5.0")
            box_x = item_left
            box_y = Inches(1.75)
            box_w = Inches(5.0)
            box_h = Inches(5.35)

            img_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
            img_box.fill.solid()
            img_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
            img_box.line.color.rgb = C_BORDER

            has_image = False
            img_path = find_image_path(item)

            if img_path:
                try:
                    im = PILImage.open(img_path)
                    im_w, im_h = im.size
                    if im_w > 0 and im_h > 0:
                        aspect = im_w / im_h

                        max_w = 4.7
                        max_h = 5.1

                        if aspect > (max_w / max_h):
                            w = max_w
                            h = max_w / aspect
                        else:
                            h = max_h
                            w = max_h * aspect

                        img_left_in = (5.0 - w) / 2
                        img_top_in = 1.75 + (5.35 - h) / 2

                        slide.shapes.add_picture(
                            img_path,
                            box_x + Inches(img_left_in),
                            Inches(img_top_in),
                            width=Inches(w),
                            height=Inches(h)
                        )
                        has_image = True
                except Exception as e:
                    print(f"Error rendering image {img_path} in PPT: {e}")

            if not has_image:
                img_tf = img_box.text_frame
                img_tf.word_wrap = True
                img_p = img_tf.paragraphs[0]
                img_p.text = f"🛋️\n\n{prod_val}\nStyle #: {style_val}"
                img_p.font.size = Pt(14)
                img_p.font.color.rgb = RGBColor(148, 163, 184)
                img_p.alignment = PP_ALIGN.CENTER

            # Right Box: Large Specs Table (Width 7.233", Spans to right border)
            wood_val = getattr(item, 'wood_type', None) or getattr(item, 'material', '—')
            finish_val = getattr(item, 'finish_color', '—')
            length_val = getattr(item, 'size_length', 0) or 0
            breadth_val = getattr(item, 'size_breadth', 0) or 0
            height_val = getattr(item, 'size_height', 0) or 0
            cbm_val = getattr(item, 'cbm', None) or getattr(item, 'total_cbm', '—')
            price_val = getattr(item, 'price_usd', None) or getattr(item, 'usd', None) or 0
            remark_val = getattr(item, 'remark', None) or getattr(item, 'remarks', None) or 'Export Quality Specification.'

            specs_rows = [
                ("Style No", str(style_val)),
                ("Product Name", str(prod_val)),
            ]
            if include_specs:
                specs_rows.extend([
                    ("Material / Wood", str(wood_val)),
                    ("Finish / Color", str(finish_val)),
                    ("Dimensions (L×B×H)", f"{length_val} × {breadth_val} × {height_val} cm"),
                    ("CBM Volume", f"{cbm_val} CBM" if cbm_val != '—' else '—'),
                ])
            if include_price:
                specs_rows.append(("Price (USD)", f"${float(price_val):.2f}" if price_val else 'Contact Quote'))

            specs_rows.append(("Remarks / Spec", str(remark_val)))

            total_rows = len(specs_rows) + 1
            table_x = box_x + box_w + Inches(0.3)
            table_w = Inches(7.233)

            table_shape = slide.shapes.add_table(
                total_rows, 2, table_x, box_y, table_w, box_h
            )
            table = table_shape.table
            table.columns[0].width = Inches(2.2)
            table.columns[1].width = Inches(5.033)

            # Header Row
            hdr_cell = table.cell(0, 0)
            table.cell(0, 1)
            hdr_cell.fill.solid()
            hdr_cell.fill.fore_color.rgb = C_WALNUT
            hdr_tf = hdr_cell.text_frame
            hdr_p = hdr_tf.paragraphs[0]
            hdr_p.text = "PRODUCT SPECIFICATIONS & DETAILS"
            hdr_p.font.size = Pt(12)
            hdr_p.font.bold = True
            hdr_p.font.color.rgb = C_WHITE
            hdr_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            hdr_cell2 = table.cell(0, 1)
            hdr_cell2.fill.solid()
            hdr_cell2.fill.fore_color.rgb = C_WALNUT

            for r_idx, (label, val) in enumerate(specs_rows, start=1):
                cell_lbl = table.cell(r_idx, 0)
                cell_val = table.cell(r_idx, 1)

                row_bg = C_ROW_ALT if r_idx % 2 == 1 else C_WHITE
                cell_lbl.fill.solid()
                cell_lbl.fill.fore_color.rgb = row_bg
                cell_val.fill.solid()
                cell_val.fill.fore_color.rgb = row_bg

                cell_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell_val.vertical_anchor = MSO_ANCHOR.MIDDLE

                p_l = cell_lbl.text_frame.paragraphs[0]
                p_l.text = f"• {label}"
                p_l.font.size = Pt(11)
                p_l.font.bold = True
                p_l.font.color.rgb = C_DARK

                p_v = cell_val.text_frame.paragraphs[0]
                p_v.text = str(val)
                p_v.font.size = Pt(11)
                if "Price" in label:
                    p_v.font.bold = True
                    p_v.font.color.rgb = C_GOLD
                else:
                    p_v.font.color.rgb = C_SLATE

        else:
            # ── 2 ITEMS PER SLIDE: SIDE-BY-SIDE LAYOUT ──
            col_width = Inches(6.0)
            gap = Inches(0.533)
            start_left = Inches(0.4)

            for col_idx, item in enumerate(chunk):
                global_idx = (page_idx - 1) * 2 + col_idx + 1
                item_left = start_left + col_idx * (col_width + gap)

                style_val = getattr(item, 'style_no', '') or getattr(item, 'sample_id', '') or f"Item #{global_idx}"
                prod_val = getattr(item, 'product_name', 'Furniture Style')

                # Item Section Title Bar
                item_title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, item_left, Inches(1.15), Inches(6.0), Inches(0.5))
                item_title_box.fill.solid()
                item_title_box.fill.fore_color.rgb = C_WALNUT
                item_title_box.line.fill.background()
                
                t_tf = item_title_box.text_frame
                t_tf.word_wrap = True
                t_p = t_tf.paragraphs[0]
                display_title = prod_val if len(prod_val) <= 32 else prod_val[:30] + "..."
                t_p.text = f" {global_idx}. {display_title} | {style_val}"
                t_p.font.size = Pt(11)
                t_p.font.bold = True
                t_p.font.color.rgb = C_WHITE
                item_title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # Left Box: Image Container
                box_x = item_left
                box_y = Inches(1.75)
                box_w = Inches(2.5)
                box_h = Inches(5.35)

                img_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
                img_box.fill.solid()
                img_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
                img_box.line.color.rgb = C_BORDER

                has_image = False
                img_path = find_image_path(item)

                if img_path:
                    try:
                        im = PILImage.open(img_path)
                        im_w, im_h = im.size
                        if im_w > 0 and im_h > 0:
                            aspect = im_w / im_h

                            max_w = 2.3
                            max_h = 5.15

                            if aspect > (max_w / max_h):
                                w = max_w
                                h = max_w / aspect
                            else:
                                h = max_h
                                w = max_h * aspect

                            img_left_in = (2.5 - w) / 2
                            img_top_in = 1.75 + (5.35 - h) / 2

                            slide.shapes.add_picture(
                                img_path,
                                item_left + Inches(img_left_in),
                                Inches(img_top_in),
                                width=Inches(w),
                                height=Inches(h)
                            )
                            has_image = True
                    except Exception as e:
                        print(f"Error rendering image {img_path} in PPT: {e}")

                if not has_image:
                    img_tf = img_box.text_frame
                    img_tf.word_wrap = True
                    img_p = img_tf.paragraphs[0]
                    img_p.text = f"🛋️\n\n{prod_val}\nStyle #: {style_val}"
                    img_p.font.size = Pt(12)
                    img_p.font.color.rgb = RGBColor(148, 163, 184)
                    img_p.alignment = PP_ALIGN.CENTER

                # Right Box: Structured Specs Table
                wood_val = getattr(item, 'wood_type', None) or getattr(item, 'material', '—')
                finish_val = getattr(item, 'finish_color', '—')
                length_val = getattr(item, 'size_length', 0) or 0
                breadth_val = getattr(item, 'size_breadth', 0) or 0
                height_val = getattr(item, 'size_height', 0) or 0
                cbm_val = getattr(item, 'cbm', None) or getattr(item, 'total_cbm', '—')
                price_val = getattr(item, 'price_usd', None) or getattr(item, 'usd', None) or 0
                remark_val = getattr(item, 'remark', None) or getattr(item, 'remarks', None) or 'Export Quality Specification.'

                specs_rows = [
                    ("Style No", str(style_val)),
                    ("Product", str(prod_val)),
                ]
                if include_specs:
                    specs_rows.extend([
                        ("Material", str(wood_val)),
                        ("Finish", str(finish_val)),
                        ("Size (L×B×H)", f"{length_val}×{breadth_val}×{height_val} cm"),
                        ("Volume", f"{cbm_val} CBM" if cbm_val != '—' else '—'),
                    ])
                if include_price:
                    specs_rows.append(("Price (USD)", f"${float(price_val):.2f}" if price_val else 'Contact Quote'))

                specs_rows.append(("Remarks", str(remark_val)))

                total_rows = len(specs_rows) + 1
                table_shape = slide.shapes.add_table(
                    total_rows, 2, item_left + Inches(2.6), box_y, Inches(3.4), box_h
                )
                table = table_shape.table
                table.columns[0].width = Inches(1.1)
                table.columns[1].width = Inches(2.3)

                # Header Row
                hdr_cell = table.cell(0, 0)
                table.cell(0, 1)
                hdr_cell.fill.solid()
                hdr_cell.fill.fore_color.rgb = C_WALNUT
                hdr_tf = hdr_cell.text_frame
                hdr_p = hdr_tf.paragraphs[0]
                hdr_p.text = "PRODUCT DETAILS"
                hdr_p.font.size = Pt(10)
                hdr_p.font.bold = True
                hdr_p.font.color.rgb = C_WHITE
                hdr_cell.vertical_anchor = MSO_ANCHOR.MIDDLE

                hdr_cell2 = table.cell(0, 1)
                hdr_cell2.fill.solid()
                hdr_cell2.fill.fore_color.rgb = C_WALNUT

                for r_idx, (label, val) in enumerate(specs_rows, start=1):
                    cell_lbl = table.cell(r_idx, 0)
                    cell_val = table.cell(r_idx, 1)

                    row_bg = C_ROW_ALT if r_idx % 2 == 1 else C_WHITE
                    cell_lbl.fill.solid()
                    cell_lbl.fill.fore_color.rgb = row_bg
                    cell_val.fill.solid()
                    cell_val.fill.fore_color.rgb = row_bg

                    cell_lbl.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell_val.vertical_anchor = MSO_ANCHOR.MIDDLE

                    p_l = cell_lbl.text_frame.paragraphs[0]
                    p_l.text = f"• {label}"
                    p_l.font.size = Pt(8.5)
                    p_l.font.bold = True
                    p_l.font.color.rgb = C_DARK

                    p_v = cell_val.text_frame.paragraphs[0]
                    p_v.text = str(val)
                    p_v.font.size = Pt(8.5)
                    if "Price" in label:
                        p_v.font.bold = True
                        p_v.font.color.rgb = C_GOLD
                    else:
                        p_v.font.color.rgb = C_SLATE


    # ── SLIDE N+1: Ending / Thank You Page ──
    end_slide = prs.slides.add_slide(blank_layout)
    
    e_bg = end_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    e_bg.fill.solid()
    e_bg.fill.fore_color.rgb = C_WALNUT
    e_bg.line.fill.background()

    # Center Card Container Box
    e_card = end_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    e_card.fill.solid()
    e_card.fill.fore_color.rgb = C_WHITE
    e_card.line.color.rgb = C_BORDER

    e_tf = e_card.text_frame
    e_tf.word_wrap = True

    ep1 = e_tf.paragraphs[0]
    ep1.text = "THANK YOU FOR YOUR INTEREST!"
    ep1.font.size = Pt(34)
    ep1.font.bold = True
    ep1.font.color.rgb = C_WALNUT
    ep1.alignment = PP_ALIGN.CENTER

    ep2 = e_tf.add_paragraph()
    ep2.text = "We look forward to partnering with you on this exclusive collection."
    ep2.font.size = Pt(16)
    ep2.font.color.rgb = C_SLATE
    ep2.alignment = PP_ALIGN.CENTER

    ep_space = e_tf.add_paragraph()
    ep_space.text = "──────────────────────────────────────────────"
    ep_space.font.size = Pt(12)
    ep_space.font.color.rgb = C_BORDER
    ep_space.alignment = PP_ALIGN.CENTER

    ep3 = e_tf.add_paragraph()
    ep3.text = "PINKCITY ENTERPRISES"
    ep3.font.size = Pt(20)
    ep3.font.bold = True
    ep3.font.color.rgb = C_DARK
    ep3.alignment = PP_ALIGN.CENTER

    ep4 = e_tf.add_paragraph()
    ep4.text = "📍 Office & Works: G-78, EPIP, Sitapura Industrial Area, Tonk Road, Jaipur-302022 Rajasthan, India."
    ep4.font.size = Pt(12)
    ep4.font.color.rgb = C_SLATE
    ep4.alignment = PP_ALIGN.CENTER

    ep5 = e_tf.add_paragraph()
    ep5.text = "📞 Tele #: +91-141-2771144 / 2770033   |   📋 GSTIN/UIN: 08ABXPS4077R1Z8   |   IEC: 1397002620   |   State Code: 08"
    ep5.font.size = Pt(12)
    ep5.font.bold = True
    ep5.font.color.rgb = C_WALNUT
    ep5.alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_brand_pptx_presentation(buyer_name="", buyer_po_numbers="", title="BRAND PRESENTATION", company_name="PINKCITY ENTERPRISES", slides_data=None):
    """
    Generates a 16:9 widescreen PowerPoint Presentation for Brand PPT / Buyer Inspection Deck.
    - Slide 1: Cover Page with Company Logo, Title, Buyer Name, Buyer PO Numbers, Date.
    - Slide 2..N: Product Collage Slides with 5-15+ images per product (tags, labels, stickers, equipment, barcode).
    - Slide N+1: Ending / Thank You Page.
    """
    if not slides_data:
        slides_data = []

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    C_WALNUT = RGBColor(139, 90, 43)      # #8B5A2B Brand Theme
    C_GOLD = RGBColor(217, 119, 6)        # #D97706 Gold Accent
    C_CREAM_BG = RGBColor(248, 246, 242)  # #F8F6F2 Light Luxury Background
    C_DARK = RGBColor(30, 41, 59)         # #1E293B Primary Text
    C_SLATE = RGBColor(71, 85, 105)       # #475569 Subtitle Text
    C_MUTED = RGBColor(100, 116, 139)     # #64748B Label Muted Text
    C_WHITE = RGBColor(255, 255, 255)
    C_BORDER = RGBColor(203, 213, 225)    # #CBD5E1 Slate Light Border

    # ── SLIDE 1: Cover Page ──
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_CREAM_BG
    bg1.line.fill.background()

    # Brand Header Bar
    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.3))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = C_WALNUT
    top_bar.line.fill.background()

    # Company Logo Insertion
    logo_path = os.path.join(settings.BASE_DIR, '..', 'Frontend', 'src', 'assets', 'Pinkcity_Logo.png')
    if not os.path.exists(logo_path):
        logo_path = r"C:\Users\User\OneDrive\Desktop\ERP Furniture\Frontend\src\assets\Pinkcity_Logo.png"

    text_left = Inches(0.4)
    if os.path.exists(logo_path):
        try:
            slide1.shapes.add_picture(logo_path, Inches(0.4), Inches(0.15), width=Inches(1.0), height=Inches(1.0))
            text_left = Inches(1.6)
        except Exception as e:
            print("Error rendering logo in Brand PPT:", e)

    tf_brand = slide1.shapes.add_textbox(text_left, Inches(0.25), Inches(11.0), Inches(0.8)).text_frame
    tf_brand.word_wrap = True
    p_brand = tf_brand.paragraphs[0]
    p_brand.text = (company_name or "PINKCITY ENTERPRISES").upper()
    p_brand.font.size = Pt(26)
    p_brand.font.bold = True
    p_brand.font.color.rgb = C_WHITE

    gold_line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.3), Inches(13.333), Inches(0.06))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = C_GOLD
    gold_line.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    
    p_title = tf_title.paragraphs[0]
    p_title.text = (title or "BRAND PRESENTATION").upper()
    p_title.font.size = Pt(34)
    p_title.font.bold = True
    p_title.font.color.rgb = C_DARK

    p_sub = tf_title.add_paragraph()
    p_sub.text = "Product Equipment, Tags & Compliance Photo Deck"
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = C_WALNUT

    card_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.6))
    card_box.fill.solid()
    card_box.fill.fore_color.rgb = C_WHITE
    card_box.line.color.rgb = C_BORDER

    tf_info = card_box.text_frame
    tf_info.word_wrap = True
    
    p_b1 = tf_info.paragraphs[0]
    p_b1.text = "PREPARED SPECIALLY FOR:"
    p_b1.font.size = Pt(11)
    p_b1.font.bold = True
    p_b1.font.color.rgb = C_MUTED

    p_b2 = tf_info.add_paragraph()
    p_b2.text = buyer_name or "Valued Buyer"
    p_b2.font.size = Pt(22)
    p_b2.font.bold = True
    p_b2.font.color.rgb = C_DARK

    if buyer_po_numbers:
        p_po = tf_info.add_paragraph()
        p_po.text = f"BUYER PO NO(S): {buyer_po_numbers}"
        p_po.font.size = Pt(15)
        p_po.font.bold = True
        p_po.font.color.rgb = C_WALNUT

    p_dt = tf_info.add_paragraph()
    p_dt.text = f"Date: {timezone.now().strftime('%d %B %Y')}   |   Total Products: {len(slides_data)}"
    p_dt.font.size = Pt(13)
    p_dt.font.color.rgb = C_SLATE

    # ── SLIDE 2..N: Product Collage Slides ──
    for slide_idx, s_data in enumerate(slides_data, start=1):
        prod_title = s_data.get('title') or s_data.get('product_name') or f"Product #{slide_idx}"
        images = s_data.get('images', [])

        slide = prs.slides.add_slide(blank_layout)

        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = RGBColor(250, 250, 249)
        s_bg.line.fill.background()

        # Top Header Banner
        h_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.3), Inches(12.533), Inches(0.65))
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = C_WALNUT
        h_bar.line.fill.background()

        h_tf = h_bar.text_frame
        h_tf.word_wrap = True
        h_p = h_tf.paragraphs[0]
        po_str = f"  |  PO: {buyer_po_numbers}" if buyer_po_numbers else ""
        h_p.text = f"  {(company_name or 'PINKCITY ENTERPRISES').upper()}   |   PRODUCT {slide_idx} OF {len(slides_data)}: {prod_title.upper()}{po_str}"
        h_p.font.size = Pt(14)
        h_p.font.bold = True
        h_p.font.color.rgb = C_WHITE
        h_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if not images:
            no_img_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(1.5))
            tf_no = no_img_box.text_frame
            p_no = tf_no.paragraphs[0]
            p_no.text = f"No photos provided for {prod_title}"
            p_no.font.size = Pt(18)
            p_no.font.color.rgb = C_MUTED
            p_no.alignment = PP_ALIGN.CENTER
            continue

        # Smart Grid Collage Placement
        n_img = len(images)
        if n_img <= 3:
            R, C = 1, n_img
        elif n_img == 4:
            R, C = 2, 2
        elif n_img <= 6:
            R, C = 2, 3
        elif n_img <= 8:
            R, C = 2, 4
        elif n_img <= 10:
            R, C = 2, 5
        elif n_img <= 12:
            R, C = 2, 6
        elif n_img <= 14:
            R, C = 2, 7
        elif n_img <= 18:
            R, C = 3, 6
        else:
            R = 3
            C = (n_img + 2) // 3

        grid_left = Inches(0.4)
        grid_top = Inches(1.1)
        grid_width = Inches(12.533)
        grid_height = Inches(6.0)

        gap_x = Inches(0.12)
        gap_y = Inches(0.12)

        cell_w = (grid_width - (C - 1) * gap_x) / C
        cell_h = (grid_height - (R - 1) * gap_y) / R

        for idx, img_input in enumerate(images):
            r = idx // C
            c = idx % C
            if r >= R:
                break # Grid full

            c_left = grid_left + c * (cell_w + gap_x)
            c_top = grid_top + r * (cell_h + gap_y)

            # Container Box Frame
            frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, c_top, cell_w, cell_h)
            frame.fill.solid()
            frame.fill.fore_color.rgb = C_WHITE
            frame.line.color.rgb = C_BORDER
            frame.line.width = Pt(1)

            try:
                pil_img = None
                img_stream = None

                if isinstance(img_input, (str, os.PathLike)) and os.path.exists(img_input):
                    pil_img = PILImage.open(img_input)
                    img_stream = img_input
                elif isinstance(img_input, bytes):
                    pil_img = PILImage.open(BytesIO(img_input))
                    img_stream = BytesIO(img_input)
                elif hasattr(img_input, 'read'):
                    img_bytes = img_input.read()
                    if hasattr(img_input, 'seek'):
                        img_input.seek(0)
                    pil_img = PILImage.open(BytesIO(img_bytes))
                    img_stream = BytesIO(img_bytes)
                elif isinstance(img_input, PILImage.Image):
                    pil_img = img_input
                    buf = BytesIO()
                    pil_img.save(buf, format='JPEG')
                    img_stream = BytesIO(buf.getvalue())
                
                if pil_img and img_stream:
                    w, h = pil_img.size
                    if w > 0 and h > 0:
                        aspect_ratio = w / h
                        pad = Inches(0.04)
                        max_w = cell_w - (2 * pad)
                        max_h = cell_h - (2 * pad)

                        if aspect_ratio > (max_w / max_h):
                            target_w = max_w
                            target_h = target_w / aspect_ratio
                        else:
                            target_h = max_h
                            target_w = target_h * aspect_ratio

                        img_l = c_left + (cell_w - target_w) / 2
                        img_t = c_top + (cell_h - target_h) / 2

                        slide.shapes.add_picture(img_stream, img_l, img_t, width=target_w, height=target_h)
            except Exception as ex:
                print(f"Error adding image {idx} to brand slide {slide_idx}: {ex}")

    # ── SLIDE N+1: Closing / Thank You ──
    end_slide = prs.slides.add_slide(blank_layout)
    e_bg = end_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    e_bg.fill.solid()
    e_bg.fill.fore_color.rgb = C_WALNUT
    e_bg.line.fill.background()

    e_card = end_slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.1))
    e_card.fill.solid()
    e_card.fill.fore_color.rgb = C_WHITE
    e_card.line.color.rgb = C_BORDER

    e_tf = e_card.text_frame
    e_tf.word_wrap = True

    ep1 = e_tf.paragraphs[0]
    ep1.text = "BRAND COMPLIANCE & QUALITY PRESENTATION"
    ep1.font.size = Pt(28)
    ep1.font.bold = True
    ep1.font.color.rgb = C_WALNUT
    ep1.alignment = PP_ALIGN.CENTER

    ep2 = e_tf.add_paragraph()
    ep2.text = f"Verified for {buyer_name or 'Valued Buyer'}" + (f" | PO #{buyer_po_numbers}" if buyer_po_numbers else "")
    ep2.font.size = Pt(16)
    ep2.font.color.rgb = C_SLATE
    ep2.alignment = PP_ALIGN.CENTER

    ep_space = e_tf.add_paragraph()
    ep_space.text = "──────────────────────────────────────────────"
    ep_space.font.size = Pt(12)
    ep_space.font.color.rgb = C_BORDER
    ep_space.alignment = PP_ALIGN.CENTER

    ep3 = e_tf.add_paragraph()
    ep3.text = (company_name or "PINKCITY ENTERPRISES").upper()
    ep3.font.size = Pt(20)
    ep3.font.bold = True
    ep3.font.color.rgb = C_DARK
    ep3.alignment = PP_ALIGN.CENTER

    ep4 = e_tf.add_paragraph()
    ep4.text = "📍 Office & Works: G-78, EPIP, Sitapura Industrial Area, Tonk Road, Jaipur-302022 Rajasthan, India."
    ep4.font.size = Pt(12)
    ep4.font.color.rgb = C_SLATE
    ep4.alignment = PP_ALIGN.CENTER

    ep5 = e_tf.add_paragraph()
    ep5.text = "📞 Tele #: +91-141-2771144 / 2770033   |   📋 GSTIN/UIN: 08ABXPS4077R1Z8   |   IEC: 1397002620"
    ep5.font.size = Pt(12)
    ep5.font.bold = True
    ep5.font.color.rgb = C_WALNUT
    ep5.alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_vendor_inspection_pptx(cover_info=None, slides_data=None):
    """
    Generates a 16:9 widescreen PowerPoint Presentation for Vendor Internal Inspection Report.
    - Slide 1: Pink city internal inspection report (Metadata & Level 1 / Level 2 AQL Tables).
    - Slide 2: Master Carton Taping & Color Coding based on Banner & DC.
    - Slide 3..N: Product Inspection Collage Slides with dynamic section headers & photo grids.
    - Slide N+1: Internal DQA Inspection Report summary slide.
    - Slide N+2: Closing / Thank You Slide.
    """
    if not cover_info:
        cover_info = {}
    if not slides_data:
        slides_data = []

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette Tokens
    C_RED_HEADER = RGBColor(235, 0, 0)     # #EB0000 Compliance Red Banner
    C_YELLOW_BANNER = RGBColor(255, 255, 0) # #FFFF00 Disclaimer Yellow
    C_DARK = RGBColor(20, 20, 20)          # Primary Text
    C_WHITE = RGBColor(255, 255, 255)
    C_GRAY_BG = RGBColor(250, 250, 250)
    C_BORDER = RGBColor(180, 180, 180)
    C_WALNUT = RGBColor(139, 90, 43)

    # ── SLIDE 1: Pink City Internal Inspection Report ──
    slide1 = prs.slides.add_slide(blank_layout)

    # Background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_WHITE
    bg1.line.fill.background()

    # Red Top Banner
    h1_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.7))
    h1_bar.fill.solid()
    h1_bar.fill.fore_color.rgb = C_RED_HEADER
    h1_bar.line.fill.background()

    h1_tf = h1_bar.text_frame
    h1_tf.word_wrap = True
    h1_p = h1_tf.paragraphs[0]
    h1_p.text = "Pinkcity - Internal Inspection Report"
    h1_p.font.size = Pt(22)
    h1_p.font.bold = True
    h1_p.font.color.rgb = C_DARK
    h1_p.alignment = PP_ALIGN.CENTER
    h1_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Left Column Metadata Fields
    left_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(4.5))
    tf_meta = left_box.text_frame
    tf_meta.word_wrap = True

    meta_items = [
        ("DQA – ", cover_info.get('dqa_name', 'Mahendra Singh')),
        ("Vendor – ", cover_info.get('vendor_name', 'Pinkcity Enterprises')),
        ("Date – ", cover_info.get('date', timezone.now().strftime('%d-%m-%Y'))),
        ("PO # ", f"{cover_info.get('po_number', '626890')} / QTY # {cover_info.get('qty', '300 / 300')}"),
        ("Ship window – ", cover_info.get('ship_window', '31 March 2026 To 10 April 2026')),
        ("Banner – ", cover_info.get('banner', 'Home Goods')),
        ("Test Report Number ", f"{cover_info.get('test_report', '(6726)041-0355 & Date - 16 Feb. 2026')}"),
        ("QEM Date - ", cover_info.get('qem_date', ''))
    ]

    for idx, (label, val) in enumerate(meta_items):
        p = tf_meta.paragraphs[0] if idx == 0 else tf_meta.add_paragraph()
        p.space_after = Pt(10)

        run_label = p.add_run()
        run_label.text = label
        run_label.font.bold = True
        run_label.font.size = Pt(14)
        run_label.font.color.rgb = C_DARK

        run_val = p.add_run()
        run_val.text = str(val)
        run_val.font.bold = True
        run_val.font.size = Pt(14)
        run_val.font.color.rgb = C_DARK

    # Level 1 Table (Right side)
    l1_title_box = slide1.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(2.8), Inches(0.3))
    tf_l1t = l1_title_box.text_frame
    p_l1t = tf_l1t.paragraphs[0]
    p_l1t.text = "LEVEL-1"
    p_l1t.font.bold = True
    p_l1t.font.size = Pt(11)
    p_l1t.font.color.rgb = C_DARK

    l1_rows = [
        ["Lot Size", "Sample Size", "Pass", "Fail"],
        ["2-8", "2", "0", "1"],
        ["9-15", "2", "0", "1"],
        ["16-25", "3", "0", "1"],
        ["51-90", "5", "0", "1"],
        ["91-150", "8", "0", "1"],
        ["151-280", "13", "0", "1"],
        ["281-500", "20", "1", "2"],
        ["501-1200", "32", "2", "3"],
        ["1201-3200", "50", "3", "4"],
        ["3201-10000", "80", "5", "6"],
        ["10001-35000", "125", "7", "8"],
        ["35000 - Above", "200", "10", "11"],
    ]

    t1_shape = slide1.shapes.add_table(len(l1_rows), 4, Inches(6.8), Inches(1.6), Inches(2.8), Inches(4.5))
    t1 = t1_shape.table
    t1.columns[0].width = Inches(1.0)
    t1.columns[1].width = Inches(0.7)
    t1.columns[2].width = Inches(0.5)
    t1.columns[3].width = Inches(0.6)

    for r_idx, row in enumerate(l1_rows):
        for c_idx, cell_text in enumerate(row):
            cell = t1.cell(r_idx, c_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(8.5)
            if r_idx == 0:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(225, 235, 225)

    # Level 2 Table (Right side)
    l2_title_box = slide1.shapes.add_textbox(Inches(9.8), Inches(1.3), Inches(3.0), Inches(0.3))
    tf_l2t = l2_title_box.text_frame
    p_l2t = tf_l2t.paragraphs[0]
    p_l2t.text = "LEVEL-2"
    p_l2t.font.bold = True
    p_l2t.font.size = Pt(11)
    p_l2t.font.color.rgb = C_DARK

    l2_rows = [
        ["Lot Size", "Sample Size", "Pass", "Fail"],
        ["2-8", "2", "0", "1"],
        ["9-15", "3", "0", "1"],
        ["16-25", "5", "0", "1"],
        ["26-50", "8", "0", "1"],
        ["51-90", "13", "0", "1"],
        ["91-150", "20", "1", "2"],
        ["151-280", "32", "2", "3"],
        ["281-500", "50", "3", "4"],
        ["501-1200", "80", "5", "6"],
        ["1201-3200", "125", "7", "8"],
        ["3201-10000", "200", "10", "11"],
        ["10001-35000", "315", "14", "15"],
        ["35001 - Above", "500", "21", "22"],
    ]

    t2_shape = slide1.shapes.add_table(len(l2_rows), 4, Inches(9.8), Inches(1.6), Inches(3.0), Inches(4.8))
    t2 = t2_shape.table
    t2.columns[0].width = Inches(1.1)
    t2.columns[1].width = Inches(0.7)
    t2.columns[2].width = Inches(0.6)
    t2.columns[3].width = Inches(0.6)

    for r_idx, row in enumerate(l2_rows):
        for c_idx, cell_text in enumerate(row):
            cell = t2.cell(r_idx, c_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(8.5)
            if r_idx == 0:
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(235, 235, 235)

    # Bottom Yellow Banner
    y_banner = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(6.5), Inches(6.3), Inches(0.5))
    y_banner.fill.solid()
    y_banner.fill.fore_color.rgb = C_YELLOW_BANNER
    y_banner.line.fill.background()

    y_tf = y_banner.text_frame
    y_p = y_tf.paragraphs[0]
    y_p.text = "Furniture - Level 2 // Non-Furniture – Level 1"
    y_p.font.bold = True
    y_p.font.size = Pt(16)
    y_p.font.color.rgb = C_DARK
    y_p.alignment = PP_ALIGN.CENTER
    y_banner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ── SLIDE 2: Master Carton Taping & Color Coding ──
    slide2 = prs.slides.add_slide(blank_layout)

    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = C_WHITE
    bg2.line.fill.background()

    # Red Top Banner
    h2_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(0.3), Inches(11.533), Inches(0.65))
    h2_bar.fill.solid()
    h2_bar.fill.fore_color.rgb = C_RED_HEADER
    h2_bar.line.fill.background()

    h2_tf = h2_bar.text_frame
    h2_p = h2_tf.paragraphs[0]
    h2_p.text = "Master Carton Taping & Color Coding based on Banner & DC"
    h2_p.font.size = Pt(20)
    h2_p.font.bold = True
    h2_p.font.color.rgb = C_DARK
    h2_p.alignment = PP_ALIGN.CENTER
    h2_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Left Side DC Dots Grid
    dots_data = [
        # (Number, Label, RGBColor)
        ("10", "DC 10", RGBColor(255, 0, 0)),        # Red
        ("50", "DC 50", RGBColor(0, 176, 240)),      # Cyan
        ("80", "DC 80", RGBColor(220, 220, 220)),    # Light Grey
        ("20", "DC 20", RGBColor(255, 192, 0)),      # Yellow
        ("60", "DC 60", RGBColor(0, 32, 96)),        # Navy
        ("30", "DC 30", RGBColor(0, 0, 0)),          # Black
        ("70", "DC 70", RGBColor(112, 48, 160)),     # Purple
        ("40", "DC 40", RGBColor(146, 208, 80)),     # Light Green
        ("90", "DC 90", RGBColor(197, 90, 17)),      # Brown
    ]

    dot_start_x = Inches(0.5)
    dot_start_y = Inches(1.3)
    col_width = Inches(2.8)
    row_height = Inches(1.1)

    for i, (num, label, color) in enumerate(dots_data):
        col = i % 3 if i < 3 else (0 if i in [3,6] else (1 if i in [4,7] else 2))
        if i in [0, 1, 2]:
            r_x = dot_start_x + (i * col_width)
            r_y = dot_start_y
        elif i in [3, 4]:
            r_x = dot_start_x + ((i - 3) * col_width)
            r_y = dot_start_y + row_height
        elif i in [5, 6]:
            r_x = dot_start_x + ((i - 5) * col_width)
            r_y = dot_start_y + (2 * row_height)
        else: # 7, 8
            r_x = dot_start_x + ((i - 7) * col_width)
            r_y = dot_start_y + (3 * row_height)

        # Circle Shape
        oval = slide2.shapes.add_shape(MSO_SHAPE.OVAL, r_x, r_y, Inches(0.85), Inches(0.85))
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background()

        ov_tf = oval.text_frame
        ov_p = ov_tf.paragraphs[0]
        ov_p.text = num
        ov_p.font.bold = True
        ov_p.font.size = Pt(16)
        ov_p.font.color.rgb = C_WHITE if color != RGBColor(220, 220, 220) and color != RGBColor(255, 192, 0) else C_DARK
        ov_p.alignment = PP_ALIGN.CENTER
        oval.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Label Text next to circle
        lbl_box = slide2.shapes.add_textbox(r_x + Inches(0.95), r_y + Inches(0.2), Inches(1.6), Inches(0.5))
        lbl_tf = lbl_box.text_frame
        lbl_p = lbl_tf.paragraphs[0]
        lbl_p.text = label
        lbl_p.font.bold = True
        lbl_p.font.size = Pt(16)
        lbl_p.font.color.rgb = C_DARK

    # Bottom Yellow Note for DC Stickers
    dc_y_banner = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(6.5), Inches(7.8), Inches(0.45))
    dc_y_banner.fill.solid()
    dc_y_banner.fill.fore_color.rgb = C_YELLOW_BANNER
    dc_y_banner.line.fill.background()

    dc_y_tf = dc_y_banner.text_frame
    dc_y_p = dc_y_tf.paragraphs[0]
    dc_y_p.text = "DC Dot Sticker Size – 2” / Placement - all 4 sides, right hand side top corner"
    dc_y_p.font.bold = True
    dc_y_p.font.size = Pt(11)
    dc_y_p.font.color.rgb = C_DARK
    dc_y_p.alignment = PP_ALIGN.CENTER
    dc_y_banner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Right Side Carton Taping Guidelines Visual Cards
    taping_bannners = [
        ("HomeGoods White color Tape", RGBColor(255, 255, 255), RGBColor(139, 90, 43)),
        ("Marshalls Black color Tape", RGBColor(0, 0, 0), RGBColor(139, 90, 43)),
        ("TJ Maxx Blue color Tape", RGBColor(0, 112, 192), RGBColor(139, 90, 43)),
    ]

    right_start_x = Inches(8.6)
    right_start_y = Inches(1.3)
    card_h = Inches(1.6)

    for idx, (t_title, tape_color, box_color) in enumerate(taping_bannners):
        curr_y = right_start_y + (idx * (card_h + Inches(0.2)))

        # Title
        tb_box = slide2.shapes.add_textbox(right_start_x, curr_y, Inches(4.3), Inches(0.35))
        tb_tf = tb_box.text_frame
        tb_p = tb_tf.paragraphs[0]
        tb_p.text = t_title
        tb_p.font.bold = True
        tb_p.font.size = Pt(13)
        tb_p.font.color.rgb = C_DARK

        # Visual Carton Diagrams (2 Boxes side by side)
        for box_i in range(2):
            bx_left = right_start_x + (box_i * Inches(2.0))
            bx_top = curr_y + Inches(0.35)

            # Carton Box Body
            c_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx_left, bx_top, Inches(1.6), Inches(1.2))
            c_box.fill.solid()
            c_box.fill.fore_color.rgb = box_color
            c_box.line.color.rgb = RGBColor(100, 60, 20)

            # Tape Line Overlay across center & edges
            tape1 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx_left + Inches(0.1), bx_top + Inches(0.45), Inches(1.4), Inches(0.25))
            tape1.fill.solid()
            tape1.fill.fore_color.rgb = tape_color
            tape1.line.color.rgb = C_DARK if tape_color == C_WHITE else tape_color

            tape2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, bx_left + Inches(0.65), bx_top + Inches(0.1), Inches(0.3), Inches(1.0))
            tape2.fill.solid()
            tape2.fill.fore_color.rgb = tape_color
            tape2.line.color.rgb = C_DARK if tape_color == C_WHITE else tape_color

    # ── SLIDES 3..N: Product Inspection Collage Slides ──
    for slide_idx, s_data in enumerate(slides_data, start=1):
        sec_title = s_data.get('title') or f"Inspection Section #{slide_idx}"
        images = s_data.get('images', [])

        slide = prs.slides.add_slide(blank_layout)

        # White background
        s_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        s_bg.fill.solid()
        s_bg.fill.fore_color.rgb = C_WHITE
        s_bg.line.fill.background()

        # Top Red Header Banner
        h_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.65))
        h_bar.fill.solid()
        h_bar.fill.fore_color.rgb = C_RED_HEADER
        h_bar.line.fill.background()

        h_tf = h_bar.text_frame
        h_tf.word_wrap = True
        h_p = h_tf.paragraphs[0]
        h_p.text = sec_title
        h_p.font.size = Pt(20)
        h_p.font.bold = True
        h_p.font.color.rgb = C_DARK
        h_p.alignment = PP_ALIGN.CENTER
        h_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if not images:
            no_img_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(1.5))
            tf_no = no_img_box.text_frame
            p_no = tf_no.paragraphs[0]
            p_no.text = f"No photos uploaded for section: {sec_title}"
            p_no.font.size = Pt(18)
            p_no.font.color.rgb = RGBColor(120, 120, 120)
            p_no.alignment = PP_ALIGN.CENTER
            continue

        num_imgs = len(images)
        top_start = Inches(1.1)
        avail_width = Inches(12.333)
        avail_height = Inches(6.0)
        margin_left = Inches(0.5)

        # Dynamic Grid Layout Logic
        if num_imgs == 1:
            rows, cols_per_row = 1, [1]
        elif num_imgs == 2:
            rows, cols_per_row = 1, [2]
        elif num_imgs == 3:
            rows, cols_per_row = 1, [3]
        elif num_imgs == 4:
            rows, cols_per_row = 2, [2, 2]
        elif num_imgs in [5, 6, 7]:
            # Matching Image 3 layout: Row 1 has 3 images, Row 2 has remaining (e.g. 4)
            top_cnt = 3
            bot_cnt = num_imgs - top_cnt
            rows, cols_per_row = 2, [top_cnt, bot_cnt]
        else: # 8 or more
            top_cnt = 4
            bot_cnt = min(num_imgs - top_cnt, 4)
            rows, cols_per_row = 2, [top_cnt, bot_cnt]

        row_h = avail_height / rows - Inches(0.15)
        img_counter = 0

        for r_idx in range(rows):
            c_count = cols_per_row[r_idx]
            cell_w = avail_width / c_count - Inches(0.15)
            r_top = top_start + (r_idx * (row_h + Inches(0.15)))

            for c_idx in range(c_count):
                if img_counter >= num_imgs:
                    break
                img_input = images[img_counter]
                img_counter += 1

                c_left = margin_left + (c_idx * (cell_w + Inches(0.15)))

                # Subdued Card Container Frame
                card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, r_top, cell_w, row_h)
                card_bg.fill.solid()
                card_bg.fill.fore_color.rgb = C_WHITE
                card_bg.line.color.rgb = C_BORDER

                try:
                    pil_img = None
                    img_stream = None

                    if isinstance(img_input, (str, os.PathLike)) and os.path.exists(img_input):
                        pil_img = PILImage.open(img_input)
                        img_stream = img_input
                    elif isinstance(img_input, bytes):
                        pil_img = PILImage.open(BytesIO(img_input))
                        img_stream = BytesIO(img_input)
                    elif hasattr(img_input, 'read'):
                        img_bytes = img_input.read()
                        if hasattr(img_input, 'seek'):
                            img_input.seek(0)
                        pil_img = PILImage.open(BytesIO(img_bytes))
                        img_stream = BytesIO(img_bytes)
                    elif isinstance(img_input, PILImage.Image):
                        pil_img = img_input
                        buf = BytesIO()
                        pil_img.save(buf, format='JPEG')
                        img_stream = BytesIO(buf.getvalue())

                    if pil_img and img_stream:
                        w, h = pil_img.size
                        if w > 0 and h > 0:
                            aspect_ratio = w / h
                            pad = Inches(0.04)
                            max_w = cell_w - (2 * pad)
                            max_h = row_h - (2 * pad)

                            if aspect_ratio > (max_w / max_h):
                                target_w = max_w
                                target_h = target_w / aspect_ratio
                            else:
                                target_h = max_h
                                target_w = target_h * aspect_ratio

                            img_l = c_left + (cell_w - target_w) / 2
                            img_t = r_top + (row_h - target_h) / 2

                            slide.shapes.add_picture(img_stream, img_l, img_t, width=target_w, height=target_h)
                except Exception as ex:
                    print(f"Error adding image {img_counter} to vendor inspection slide: {ex}")

    # ── SLIDE N+1: Internal DQA Inspection Report ──
    sum_slide = prs.slides.add_slide(blank_layout)
    s_bg2 = sum_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    s_bg2.fill.solid()
    s_bg2.fill.fore_color.rgb = C_WHITE
    s_bg2.line.fill.background()

    # Red Top Banner
    sh_bar = sum_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.65))
    sh_bar.fill.solid()
    sh_bar.fill.fore_color.rgb = C_RED_HEADER
    sh_bar.line.fill.background()

    sh_tf = sh_bar.text_frame
    sh_p = sh_tf.paragraphs[0]
    sh_p.text = "Internal DQA Inspection Report"
    sh_p.font.size = Pt(20)
    sh_p.font.bold = True
    sh_p.font.color.rgb = C_DARK
    sh_p.alignment = PP_ALIGN.CENTER
    sh_bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    dqa_img_input = cover_info.get('dqa_report_image')
    if dqa_img_input:
        try:
            pil_img = None
            img_stream = None

            if isinstance(dqa_img_input, (str, os.PathLike)) and os.path.exists(dqa_img_input):
                pil_img = PILImage.open(dqa_img_input)
                img_stream = dqa_img_input
            elif isinstance(dqa_img_input, bytes):
                pil_img = PILImage.open(BytesIO(dqa_img_input))
                img_stream = BytesIO(dqa_img_input)
            elif hasattr(dqa_img_input, 'read'):
                img_bytes = dqa_img_input.read()
                if hasattr(dqa_img_input, 'seek'):
                    dqa_img_input.seek(0)
                pil_img = PILImage.open(BytesIO(img_bytes))
                img_stream = BytesIO(img_bytes)

            if pil_img and img_stream:
                w, h = pil_img.size
                if w > 0 and h > 0:
                    aspect_ratio = w / h
                    max_w = Inches(10.5)
                    max_h = Inches(6.0)

                    if aspect_ratio > (max_w / max_h):
                        target_w = max_w
                        target_h = target_w / aspect_ratio
                    else:
                        target_h = max_h
                        target_w = target_h * aspect_ratio

                    img_l = Inches(0.5) + (Inches(12.333) - target_w) / 2
                    img_t = Inches(1.1) + (Inches(6.0) - target_h) / 2

                    # Outer frame box
                    frame = sum_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_l - Inches(0.04), img_t - Inches(0.04), target_w + Inches(0.08), target_h + Inches(0.08))
                    frame.fill.solid()
                    frame.fill.fore_color.rgb = C_WHITE
                    frame.line.color.rgb = C_BORDER

                    sum_slide.shapes.add_picture(img_stream, img_l, img_t, width=target_w, height=target_h)
        except Exception as ex:
            print(f"Error adding DQA report image to slide: {ex}")
    else:
        no_img_box = sum_slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(1.5))
        tf_no = no_img_box.text_frame
        p_no = tf_no.paragraphs[0]
        p_no.text = "No Internal DQA Inspection Report Image Uploaded"
        p_no.font.size = Pt(18)
        p_no.font.color.rgb = RGBColor(120, 120, 120)
        p_no.alignment = PP_ALIGN.CENTER

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()




